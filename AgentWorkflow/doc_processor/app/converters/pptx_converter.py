# services/doc-processor/app/converters/pptx_converter.py
# Converts a PPTX file into one PNG image per slide.
# Strategy: convert PPTX → PDF first (via LibreOffice), then PDF → PNG.
# Fallback: render slides directly using python-pptx + Pillow if LibreOffice unavailable.

import os
import subprocess
import tempfile
from pathlib import Path
from typing import List


def convert_pptx_to_images(
    pptx_path: str,
    output_dir: str,
    dpi: int = 150,
) -> List[str]:
    """
    Convert each slide of a PPTX into a PNG image.

    Strategy:
    1. Try LibreOffice headless (best quality, preserves fonts/layout)
    2. Fall back to python-pptx direct rendering (basic shapes/text only)

    Args:
        pptx_path:  Path to the source PPTX file
        output_dir: Directory to save output PNG images
        dpi:        Resolution for rendering

    Returns:
        List of absolute paths to generated images, ordered by slide number
    """
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    # Try LibreOffice first
    if _libreoffice_available():
        return _convert_via_libreoffice(pptx_path, str(output_path), dpi)
    else:
        return _convert_via_python_pptx(pptx_path, str(output_path))


def _libreoffice_available() -> bool:
    """Check if LibreOffice headless is installed."""
    for cmd in ["libreoffice", "soffice"]:
        try:
            result = subprocess.run(
                [cmd, "--version"],
                capture_output=True,
                timeout=5,
            )
            if result.returncode == 0:
                return True
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return False


def _convert_via_libreoffice(pptx_path: str, output_dir: str, dpi: int) -> List[str]:
    """Convert PPTX → PDF via LibreOffice, then PDF → PNG."""
    from app.converters.pdf_converter import convert_pdf_to_images

    with tempfile.TemporaryDirectory() as tmp_dir:
        # Step 1: PPTX → PDF
        lo_cmd = next(
            (cmd for cmd in ["libreoffice", "soffice"] if _cmd_exists(cmd)),
            None,
        )
        if not lo_cmd:
            raise RuntimeError("LibreOffice not found")

        result = subprocess.run(
            [
                lo_cmd,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                pptx_path,
            ],
            capture_output=True,
            timeout=120,
        )

        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed: {result.stderr.decode()}"
            )

        # Find generated PDF
        pdf_files = list(Path(tmp_dir).glob("*.pdf"))
        if not pdf_files:
            raise RuntimeError("LibreOffice did not produce a PDF file")

        # Step 2: PDF → PNG
        return convert_pdf_to_images(
            pdf_path=str(pdf_files[0]),
            output_dir=output_dir,
            dpi=dpi,
        )


def _convert_via_python_pptx(pptx_path: str, output_dir: str) -> List[str]:
    """
    Fallback: render slides using python-pptx.
    Captures text and basic shapes — complex graphics may not render perfectly.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        raise RuntimeError(
            "python-pptx and Pillow are required for fallback rendering. "
            "Install with: pip install python-pptx Pillow"
        )

    prs = Presentation(pptx_path)
    image_paths = []

    # Standard slide dimensions (1280x720 px at 96 dpi)
    WIDTH, HEIGHT = 1280, 720
    BG_COLOR = (255, 255, 255)
    TEXT_COLOR = (30, 30, 30)
    TITLE_COLOR = (30, 30, 80)

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (WIDTH, HEIGHT), BG_COLOR)
        draw = ImageDraw.Draw(img)

        # Try to load a font, fall back to default
        try:
            title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
            body_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 22)
        except (IOError, OSError):
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()

        y_offset = 40

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para_i, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue

                # Heuristic: first text block on slide is likely the title
                is_title = (para_i == 0 and shape == slide.shapes[0])
                font = title_font if is_title else body_font
                color = TITLE_COLOR if is_title else TEXT_COLOR

                # Word-wrap text
                words = text.split()
                line = ""
                for word in words:
                    test_line = f"{line} {word}".strip()
                    bbox = draw.textbbox((0, 0), test_line, font=font)
                    if bbox[2] > WIDTH - 80:
                        draw.text((40, y_offset), line, font=font, fill=color)
                        y_offset += bbox[3] - bbox[1] + 8
                        line = word
                    else:
                        line = test_line

                if line:
                    bbox = draw.textbbox((0, 0), line, font=font)
                    draw.text((40, y_offset), line, font=font, fill=color)
                    y_offset += bbox[3] - bbox[1] + 8

                y_offset += 10  # paragraph spacing

                if y_offset > HEIGHT - 40:
                    break

        filename = f"slide_{i + 1:03d}.png"
        filepath = Path(output_dir) / filename
        img.save(str(filepath), "PNG")
        image_paths.append(str(filepath))

    return image_paths


def _cmd_exists(cmd: str) -> bool:
    try:
        subprocess.run([cmd, "--version"], capture_output=True, timeout=3)
        return True
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return False
